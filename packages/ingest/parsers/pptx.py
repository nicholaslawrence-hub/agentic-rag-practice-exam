"""Parse PPTX files into structured slide documents, with optional image rendering."""

from __future__ import annotations

import hashlib
import re
import shutil
import subprocess
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt


@dataclass
class SlideDoc:
    source: str
    slide_num: int
    title: str
    body: str
    notes: str
    week: int | None
    topic: str | None
    slide_image_url: str | None = None  # set after image rendering

    @property
    def text(self) -> str:
        parts = []
        if self.title:
            parts.append(f"[Slide {self.slide_num}: {self.title}]")
        if self.body:
            parts.append(self.body)
        if self.notes:
            parts.append(f"[Notes] {self.notes}")
        return "\n".join(parts)

    @property
    def chunk_id(self) -> str:
        h = hashlib.sha256(self.text.encode()).hexdigest()[:12]
        return f"slide::{Path(self.source).stem}::{self.slide_num}::{h}"


def _extract_text(shape) -> str:
    chunks: list[str] = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            line = " ".join(run.text for run in para.runs if run.text.strip())
            if line:
                chunks.append(line)
    if shape.shape_type == 19:  # TABLE
        for row in shape.table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                chunks.append(row_text)
    return "\n".join(chunks)


def _slide_title(slide) -> str:
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        return slide.shapes.title.text_frame.text.strip()
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size >= Pt(20):
                        if run.text.strip():
                            return run.text.strip()
    return ""


def _slide_notes(slide) -> str:
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def _find_libreoffice() -> str | None:
    """Locate the LibreOffice binary (cross-platform)."""
    for candidate in ["libreoffice", "soffice"]:
        found = shutil.which(candidate)
        if found:
            return found
    # Common Windows install paths
    for win_path in [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]:
        if Path(win_path).exists():
            return win_path
    return None


def render_pptx_to_images(
    pptx_path: Path,
    course_id: str,
    *,
    data_root: Path = Path("data"),
    dpi: int = 150,
) -> dict[int, str]:
    """
    Render each slide of a PPTX to a PNG using LibreOffice + pdf2image.

    Returns a dict mapping slide_num → relative API URL path, e.g.:
      { 1: "/static/mcb102/slide_images/Lecture_03/slide_001.png", ... }

    Returns {} if LibreOffice or pdf2image is not available (graceful degradation).
    """
    lo = _find_libreoffice()
    if not lo:
        return {}

    try:
        from pdf2image import convert_from_path  # type: ignore
    except ImportError:
        return {}

    # Output directory for PNGs
    img_dir = data_root / course_id / "slide_images" / pptx_path.stem
    img_dir.mkdir(parents=True, exist_ok=True)

    # Convert PPTX → PDF via LibreOffice
    pdf_path = img_dir / f"{pptx_path.stem}.pdf"
    try:
        subprocess.run(
            [lo, "--headless", "--convert-to", "pdf", "--outdir", str(img_dir), str(pptx_path)],
            check=True,
            capture_output=True,
            timeout=120,
        )
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired):
        return {}

    if not pdf_path.exists():
        return {}

    # PDF → per-slide PNGs
    try:
        images = convert_from_path(pdf_path, dpi=dpi)
    except Exception:
        return {}

    result: dict[int, str] = {}
    for i, img in enumerate(images, start=1):
        img_file = img_dir / f"slide_{i:03d}.png"
        img.save(img_file, "PNG")
        # Return URL relative to FastAPI static mount
        result[i] = f"/static/{course_id}/slide_images/{pptx_path.stem}/slide_{i:03d}.png"

    return result


def parse_pptx(
    path: Path,
    *,
    week: int | None = None,
    topic: str | None = None,
    slide_images: dict[int, str] | None = None,
) -> list[SlideDoc]:
    """
    Parse a PPTX into SlideDoc objects.

    slide_images: optional dict from render_pptx_to_images, maps slide_num → image URL.
    """
    prs = Presentation(str(path))
    docs: list[SlideDoc] = []
    source = path.name
    images = slide_images or {}

    for i, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide)
        body_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                txt = _extract_text(shape)
                if txt:
                    body_parts.append(txt)
            elif shape.shape_type == 19:
                txt = _extract_text(shape)
                if txt:
                    body_parts.append(txt)
        body = "\n".join(body_parts)
        notes = _slide_notes(slide)

        if not (title or body or notes):
            continue

        docs.append(SlideDoc(
            source=source,
            slide_num=i,
            title=title,
            body=body,
            notes=notes,
            week=week,
            topic=topic,
            slide_image_url=images.get(i),
        ))

    return docs
