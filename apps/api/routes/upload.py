"""Document upload - extract text from student files and store per thread."""

from __future__ import annotations

import io
import uuid
from pathlib import Path

from fastapi import APIRouter, Depends, HTTPException, Query, UploadFile
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from apps.api.auth import get_current_user
from apps.api.db.models import Thread, UserDocument
from apps.api.db.session import get_db

router = APIRouter(prefix="/upload", tags=["upload"])

_ALLOWED_EXT = {".pdf", ".pptx", ".docx", ".txt"}
_MAX_BYTES = 20 * 1024 * 1024  # 20 MB


def _extract_text(filename: str, content: bytes) -> str:
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        try:
            from pypdf import PdfReader  # type: ignore

            reader = PdfReader(io.BytesIO(content))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception:
            return ""
    if ext == ".docx":
        try:
            import docx  # type: ignore

            doc = docx.Document(io.BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception:
            return ""
    if ext in (".pptx", ".ppt"):
        try:
            from pptx import Presentation  # type: ignore

            prs = Presentation(io.BytesIO(content))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        texts.append(shape.text_frame.text)
            return "\n".join(texts)
        except Exception:
            return ""
    if ext == ".txt":
        return content.decode("utf-8", errors="ignore")
    return ""


@router.post("")
async def upload_document(
    file: UploadFile,
    thread_id: str | None = Query(default=None),
    course: str = Query(default=""),
    user=Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Upload a document (PDF/PPTX/DOCX/TXT).  Text is extracted and stored per thread."""
    ext = Path(file.filename or "").suffix.lower()
    if ext not in _ALLOWED_EXT:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type. Allowed: {', '.join(sorted(_ALLOWED_EXT))}",
        )

    content = await file.read()
    if len(content) > _MAX_BYTES:
        raise HTTPException(status_code=413, detail="File too large (max 20 MB)")

    # Resolve or create thread
    if thread_id:
        result = await db.execute(
            select(Thread).where(Thread.id == thread_id, Thread.user_id == user.id)
        )
        thread = result.scalar_one_or_none()
        if not thread:
            raise HTTPException(status_code=404, detail="Thread not found")
    else:
        if not course:
            raise HTTPException(
                status_code=400, detail="course is required when creating a new thread"
            )
        thread = Thread(id=uuid.uuid4(), user_id=user.id, course=course, title="New chat")
        db.add(thread)
        await db.flush()

    extracted = _extract_text(file.filename or "", content)[:60_000]
    doc = UserDocument(
        id=uuid.uuid4(),
        thread_id=thread.id,
        user_id=user.id,
        filename=file.filename or "upload",
        extracted_text=extracted,
    )
    db.add(doc)
    await db.commit()

    return {
        "doc_id": str(doc.id),
        "thread_id": str(thread.id),
        "filename": file.filename,
        "chars": len(extracted),
    }


@router.get("/documents")
async def list_documents(
    thread_id: str = Query(),
    user=Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    result = await db.execute(
        select(UserDocument)
        .where(UserDocument.thread_id == thread_id, UserDocument.user_id == user.id)
        .order_by(UserDocument.created_at)
    )
    docs = result.scalars().all()
    return [
        {"doc_id": str(d.id), "filename": d.filename, "chars": len(d.extracted_text)}
        for d in docs
    ]